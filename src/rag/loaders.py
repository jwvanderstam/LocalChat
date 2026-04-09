"""
Document Loaders
================

Handles loading of various document formats (text, PDF, DOCX, images).
Extracted from DocumentProcessor for modularity.
"""

import base64
import os
import re
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple, Union

from .. import config
from ..ollama_client import ollama_client
from ..utils.logging_config import get_logger

logger = get_logger(__name__)

# Document loaders with proper initialization
PyPDF2 = None
PDF_AVAILABLE = False
try:
    import PyPDF2 as pypdf2_lib
    PyPDF2 = pypdf2_lib
    PDF_AVAILABLE = True
    logger.debug("PyPDF2 available for PDF processing")
except ImportError:
    logger.warning("PyPDF2 not available - PDF support disabled")

Document = None
DOCX_AVAILABLE = False
try:
    from docx import Document as DocxDocument
    Document = DocxDocument
    DOCX_AVAILABLE = True
    logger.debug("python-docx available for DOCX processing")
except ImportError:
    logger.warning("python-docx not available - DOCX support disabled")

Presentation = None
PPTX_AVAILABLE = False
try:
    from pptx import Presentation as PptxPresentation
    Presentation = PptxPresentation
    PPTX_AVAILABLE = True
    logger.debug("python-pptx available for PPTX processing")
except ImportError:
    logger.debug("python-pptx not installed - PPTX support disabled")

# Try to import monitoring - graceful degradation if not available
try:
    from ..monitoring import counted, timed
except ImportError:
    def timed(_metric_name):  # noqa: E306
        return lambda func: func
    def counted(_metric_name, _labels=None):  # noqa: E306
        return lambda func: func

_PDF_NOT_INSTALLED = "PyPDF2 not installed"
_DOCX_NOT_INSTALLED = "python-docx not installed"
_PPTX_NOT_INSTALLED = "python-pptx not installed — run: pip install python-pptx"
_DOCX_LOAD_ERROR = "Error loading DOCX: "
_EXTRACTOR_PDFPLUMBER = "pdfplumber"
_EXTRACTOR_PYPDF2 = "PyPDF2"


class DocumentLoaderMixin:
    """
    Mixin providing document loading methods for DocumentProcessor.

    Supports text, PDF (with table extraction), DOCX, and image files.
    """

    def load_text_file(self, file_path: str) -> tuple[bool, str]:
        """
        Load a text file.

        Args:
            file_path: Path to text file

        Returns:
            Tuple of (success: bool, content_or_error: str)
        """
        try:
            logger.debug(f"Loading text file: {file_path}")
            with open(file_path, encoding='utf-8') as f:
                content = f.read()
            logger.debug(f"Loaded {len(content)} characters from text file")
            return True, content
        except Exception as e:
            logger.error(f"Error loading text file: {e}", exc_info=True)
            return False, str(e)

    def _try_pdfplumber_import(self):
        """Try to import pdfplumber, return module or None."""
        try:
            import pdfplumber as pdf_lib
            return pdf_lib
        except ImportError:
            return None

    def _format_table_rows(self, table: list) -> str:
        """Format a pdfplumber table (list of rows) into pipe-delimited text."""
        lines = ""
        for row in table:
            if not row:
                continue
            row_text = " | ".join(str(cell).strip() if cell else "" for cell in row)
            if row_text.strip():
                lines += row_text + "\n"
        return lines

    def _extract_pdfplumber_table_text(self, page, page_num: int) -> str:
        """Extract and format table text from a single pdfplumber page."""
        try:
            tables = page.extract_tables()
            if not tables:
                return ""
            table_text = ""
            for table_idx, table in enumerate(tables, 1):
                if not table:
                    continue
                table_text += f"\n[Table {table_idx} on page {page_num}]\n"
                table_text += self._format_table_rows(table)
                table_text += "\n"
            return table_text
        except Exception as e:
            logger.error(f"  Page {page_num}: Error extracting tables: {e}")
            return ""

    def _extract_pdfplumber_text(self, pdfplumber_module, file_path: str) -> str:
        """Extract full concatenated text from a PDF using pdfplumber."""
        text = ""
        with pdfplumber_module.open(file_path) as pdf:
            num_pages = len(pdf.pages)
            logger.info(f"PDF has {num_pages} pages (using pdfplumber)")
            for page_num, page in enumerate(pdf.pages, 1):
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
                else:
                    logger.warning(f"  Page {page_num}: no text extracted")
                text += self._extract_pdfplumber_table_text(page, page_num)
        logger.info(f"pdfplumber extraction: {len(text):,} chars from {num_pages} pages")
        if len(text) < 100:
            raise ValueError("Insufficient text extracted with pdfplumber")
        return text

    def _extract_pypdf2_text(self, file_path: str) -> str:
        """Extract full concatenated text from a PDF using PyPDF2."""
        text = ""
        with open(file_path, 'rb') as f:
            pdf_reader = PyPDF2.PdfReader(f)
            num_pages = len(pdf_reader.pages)
            logger.info(f"PDF has {num_pages} pages (using PyPDF2)")
            for page_num, page in enumerate(pdf_reader.pages, 1):
                try:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
                    else:
                        logger.warning(f"  Page {page_num}: no text extracted")
                except Exception as e:
                    logger.error(f"  Page {page_num}: Error extracting text: {e}")
        logger.info(f"PyPDF2 extraction complete: {len(text):,} characters")
        return text

    def _load_pages_pdfplumber(self, pdfplumber_module, file_path: str) -> list[dict[str, Any]]:
        """Load per-page data with metadata using pdfplumber."""
        pages_data = []
        with pdfplumber_module.open(file_path) as pdf:
            num_pages = len(pdf.pages)
            logger.info(f"PDF has {num_pages} pages")
            for page_num, page in enumerate(pdf.pages, 1):
                page_text = (page.extract_text() or "") + self._extract_pdfplumber_table_text(page, page_num)
                if page_text.strip():
                    pages_data.append({
                        'page_number': page_num,
                        'text': page_text,
                        'section_title': self._extract_section_title(page_text)
                    })
                else:
                    logger.warning(f"  Page {page_num}: No text extracted")
        return pages_data

    def _load_pages_pypdf2(self, file_path: str) -> list[dict[str, Any]]:
        """Load per-page data with metadata using PyPDF2."""
        pages_data = []
        with open(file_path, 'rb') as f:
            pdf_reader = PyPDF2.PdfReader(f)
            num_pages = len(pdf_reader.pages)
            logger.info(f"PDF has {num_pages} pages (PyPDF2)")
            for page_num, page in enumerate(pdf_reader.pages, 1):
                page_text = page.extract_text() or ""
                if page_text.strip():
                    pages_data.append({
                        'page_number': page_num,
                        'text': page_text,
                        'section_title': self._extract_section_title(page_text)
                    })
        return pages_data

    def load_pdf_file(self, file_path: str) -> tuple[bool, str]:
        """
        Load a PDF file with enhanced table extraction and improved text extraction.

        Uses pdfplumber for better table extraction if available,
        falls back to PyPDF2 for basic text extraction.

        Args:
            file_path: Path to PDF file

        Returns:
            Tuple of (success: bool, content_or_error: str)
        """
        if not PDF_AVAILABLE:
            logger.error(_PDF_NOT_INSTALLED)
            return False, _PDF_NOT_INSTALLED

        try:
            logger.info(f"Loading PDF file: {file_path}")
            file_size = os.path.getsize(file_path)
            logger.debug(f"PDF file size: {file_size:,} bytes")

            pdfplumber = self._try_pdfplumber_import()
            text = ""
            extraction_method = "unknown"

            if pdfplumber is not None:
                try:
                    text = self._extract_pdfplumber_text(pdfplumber, file_path)
                    extraction_method = _EXTRACTOR_PDFPLUMBER
                except Exception as plumber_error:
                    logger.warning(f"{_EXTRACTOR_PDFPLUMBER} extraction failed: {plumber_error}, falling back to {_EXTRACTOR_PYPDF2}")
                    text = ""

            if not text:
                text = self._extract_pypdf2_text(file_path)
                extraction_method = _EXTRACTOR_PYPDF2

            if not text.strip():
                error_msg = (
                    f"PDF extraction yielded no text using {extraction_method} "
                    "(file may be image-based, password-protected, or corrupt)."
                )
                logger.error(f"{error_msg} File: {file_path} ({file_size:,} bytes)")
                return False, error_msg

            if len(text) < 100:
                logger.warning(f"PDF extraction yielded very little text: {len(text)} characters")

            logger.info(f"PDF extraction successful using {extraction_method}: {len(text):,} characters extracted")
            return True, text

        except Exception as e:
            logger.error(f"Error loading PDF: {e}", exc_info=True)
            return False, str(e)

    def _load_pdf_with_pages(
        self, file_path: str
    ) -> tuple[bool, list[dict[str, Any]] | str]:
        """
        Load PDF with page-by-page tracking for enhanced citations.

        Args:
            file_path: Path to PDF file

        Returns:
            Tuple of (success: bool, pages_data_or_error)
        """
        if not PDF_AVAILABLE:
            logger.error(_PDF_NOT_INSTALLED)
            return False, _PDF_NOT_INSTALLED

        try:
            logger.info(f"Loading PDF with page tracking: {file_path}")
            pdfplumber = self._try_pdfplumber_import()

            if pdfplumber is not None:
                pages_data = self._load_pages_pdfplumber(pdfplumber, file_path)
            else:
                logger.warning("pdfplumber not available, using PyPDF2")
                pages_data = self._load_pages_pypdf2(file_path)

            logger.info(f"Extracted {len(pages_data)} pages with metadata")

            if not pages_data:
                return False, "No pages with text extracted from PDF"

            return True, pages_data

        except Exception as e:
            logger.error(f"Error loading PDF with pages: {e}", exc_info=True)
            return False, str(e)

    def _validate_docx_file(self, file_path: str) -> tuple[bool, str]:
        """Validate DOCX file exists and is non-empty. Returns (ok, error_msg)."""
        if not os.path.exists(file_path):
            return False, f"File not found: {file_path}"
        if os.path.getsize(file_path) == 0:
            return False, "File is empty (0 bytes)"
        return True, ""

    def _extract_docx_text(self, doc) -> str:
        """Extract combined text from DOCX paragraphs and table cells."""
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        cells = [
            cell.text.strip()
            for table in doc.tables
            for row in table.rows
            for cell in row.cells
            if cell.text.strip()
        ]
        return "\n".join(paragraphs + cells)

    def load_docx_file(self, file_path: str) -> tuple[bool, str]:
        """
        Load a DOCX file with enhanced error handling.

        Extracts text from both paragraphs and tables in the document.

        Args:
            file_path: Path to DOCX file

        Returns:
            Tuple of (success: bool, content_or_error: str)
        """
        if not DOCX_AVAILABLE:
            logger.error(_DOCX_NOT_INSTALLED)
            return False, _DOCX_NOT_INSTALLED

        try:
            logger.debug(f"Loading DOCX file: {file_path}")
            ok, err = self._validate_docx_file(file_path)
            if not ok:
                logger.error(err)
                return False, err

            try:
                doc = Document(file_path)
            except Exception as doc_error:
                error_msg = f"Failed to open DOCX file: {doc_error}"
                logger.error(error_msg, exc_info=True)
                return False, f"{error_msg}. File might be corrupted or password-protected."

            text = self._extract_docx_text(doc)
            logger.debug(f"DOCX extracted {len(text)} characters")

            if len(text.strip()) < 10:
                error_msg = f"Document appears to be empty or has very little text ({len(text)} chars)"
                logger.warning(error_msg)
                return False, f"{error_msg}. Check if document has actual content."

            return True, text

        except Exception as e:
            logger.error(f"{_DOCX_LOAD_ERROR}{e}", exc_info=True)
            return False, f"{_DOCX_LOAD_ERROR}{e}"

    def _line_looks_like_title(self, line: str) -> str | None:
        """Return title text if line looks like a section header, else None."""
        if len(line) >= 100:
            return None
        if line.endswith(':'):
            return line.rstrip(':')
        words = line.split()
        if len(words) >= 2 and (line.istitle() or line.isupper()):
            return line
        return None

    def _extract_section_title(self, page_text: str) -> str | None:
        """
        Extract likely section title from page start.

        Args:
            page_text: Text content of a page

        Returns:
            Section title string or None if no title found
        """
        if not page_text:
            return None
        for line in page_text.strip().split('\n')[:5]:
            line = line.strip()
            if not line or len(line) < 3 or re.match(r'^\d+\.', line):
                continue
            title = self._line_looks_like_title(line)
            if title is not None:
                return title
        return None

    def _process_pptx_slide(self, slide: Any) -> tuple[list[str], str | None]:
        """Extract text parts and title from a single PowerPoint slide."""
        title: str | None = None
        parts: list[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text or shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                continue
            if (title is None and hasattr(shape, 'placeholder_format')
                    and shape.placeholder_format is not None
                    and shape.placeholder_format.idx == 0):
                title = text
            parts.append(text)
        return parts, title

    def load_pptx_file(self, file_path: str) -> tuple[bool, list[dict[str, Any]] | str]:
        """
        Load a PowerPoint file and return a list of slide dicts.

        Each slide dict has keys: slide_number (int), text (str), title (str | None).

        Returns:
            Tuple of (success: bool, slides_or_error)
        """
        if not PPTX_AVAILABLE:
            logger.error(_PPTX_NOT_INSTALLED)
            return False, _PPTX_NOT_INSTALLED

        try:
            logger.info(f"Loading PPTX file: {file_path}")
            prs = Presentation(file_path)
            slides = []
            for slide_num, slide in enumerate(prs.slides, 1):
                parts, title = self._process_pptx_slide(slide)
                combined = "\n".join(parts)
                if combined.strip():
                    slides.append({
                        'slide_number': slide_num,
                        'text': combined,
                        'title': title,
                    })
            logger.info(f"PPTX: loaded {len(slides)} slides from {file_path}")
            if not slides:
                return False, "No text found in presentation"
            return True, slides
        except Exception as e:
            logger.error(f"Error loading PPTX: {e}", exc_info=True)
            return False, str(e)

    @staticmethod
    def _decode_email_part(part: Any) -> str | None:
        """Decode a text/plain MIME part to a string, or return None."""
        if part.get_content_type() != 'text/plain':
            return None
        charset = part.get_content_charset() or 'utf-8'
        payload = part.get_payload(decode=True)
        if not payload:
            return None
        try:
            return payload.decode(charset, errors='replace')
        except (LookupError, UnicodeDecodeError):
            return payload.decode('utf-8', errors='replace')

    def load_eml_file(self, file_path: str) -> tuple[bool, str]:
        """
        Load an .eml email file and return its plain-text body.

        Uses the stdlib ``email`` module — no extra dependency.

        Returns:
            Tuple of (success: bool, content_or_error: str)
        """
        import email as _email_mod
        try:
            logger.info(f"Loading EML file: {file_path}")
            with open(file_path, 'rb') as f:
                msg = _email_mod.message_from_bytes(f.read())

            parts: list[str] = []
            for header in ('Subject', 'From', 'Date'):
                value = msg.get(header, '')
                if value:
                    parts.append(f"{header}: {value}")

            for part in msg.walk():
                decoded = self._decode_email_part(part)
                if decoded:
                    parts.append(decoded)

            content = "\n\n".join(parts).strip()
            if not content:
                return False, "No text content found in email"
            logger.info(f"EML loaded: {len(content)} chars from {file_path}")
            return True, content
        except Exception as e:
            logger.error(f"Error loading EML: {e}", exc_info=True)
            return False, str(e)

    def load_image_file(self, file_path: str) -> tuple[bool, str]:
        """
        Load an image file by generating a text description via a vision model.

        Args:
            file_path: Path to the image file (.png, .jpg, .jpeg, .gif, .webp)

        Returns:
            Tuple of (success: bool, description_or_error: str)
        """
        try:
            logger.info(f"Loading image file: {file_path}")

            with open(file_path, 'rb') as f:
                image_data = f.read()

            image_b64 = base64.b64encode(image_data).decode('utf-8')

            vision_model = ollama_client.get_vision_model()
            if not vision_model:
                return False, "No vision model available for image processing"

            success, description = ollama_client.describe_image(vision_model, image_b64)
            if not success:
                return False, f"Failed to describe image: {description}"

            filename = os.path.basename(file_path)
            content = f"Image: {filename}\n\n{description}"
            logger.info(f"Image loaded: {len(content)} characters from {filename}")
            return True, content

        except Exception as e:
            logger.error(f"Error loading image file: {e}", exc_info=True)
            return False, str(e)

    @timed('rag.load_document')
    @counted('rag.document_loads')
    def load_document(self, file_path: str) -> tuple[bool, str]:
        """
        Load a document based on its extension.

        Args:
            file_path: Path to document file

        Returns:
            Tuple of (success: bool, content_or_error: str)
        """
        ext = Path(file_path).suffix.lower()
        logger.debug(f"Loading document with extension: {ext}")

        if ext in ('.txt', '.md', '.py', '.js', '.ts'):
            return self.load_text_file(file_path)
        elif ext == '.pdf':
            return self.load_pdf_file(file_path)
        elif ext == '.docx':
            return self.load_docx_file(file_path)
        elif ext == '.pptx':
            # PPTX returns slide list, not plain text — callers must use load_pptx_file directly
            return self.load_pptx_file(file_path)  # type: ignore[return-value]
        elif ext == '.eml':
            return self.load_eml_file(file_path)
        elif ext in config.SUPPORTED_IMAGE_EXTENSIONS:
            return self.load_image_file(file_path)
        else:
            error_msg = f"Unsupported file type: {ext}"
            logger.error(error_msg)
            return False, error_msg
